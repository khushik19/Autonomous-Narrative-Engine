# visual_agent.py — FINAL VERSION
#
# Strategy:
#   1. Open template → keep ONLY background images/shapes
#   2. DELETE every text box from template (they are all broken/off-screen)
#   3. ADD our own clean textboxes at correct positions
#   4. This gives us: template background + clean readable content
#
# Place in: narrativa_project/backend/agents/visual_agent.py

import os
import json
import copy
import re
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── Font Themes ───────────────────────
FONT_THEMES = {
    "modern":    {"title": "Segoe UI Black", "body": "Segoe UI", "accent": "Impact"},
    "corporate": {"title": "Arial Black",    "body": "Calibri",  "accent": "Georgia"},
    "luxury":    {"title": "Georgia",        "body": "Garamond", "accent": "Palatino Linotype"},
    "impact":    {"title": "Impact",         "body": "Segoe UI", "accent": "Arial Black"},
}


# ─────────────────────────────────────
# MODEL LOADING (layout model only)
# ─────────────────────────────────────

layout_model     = None
layout_tokenizer = None

def load_models():
    global layout_model, layout_tokenizer
    if layout_model is not None:
        return
    try:
        from unsloth import FastLanguageModel
        import torch

        base_dir    = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        layout_path = os.path.join(base_dir, "models", "layout_model_final")

        if not os.path.exists(layout_path):
            print(f"  Model not found at {layout_path} — using rules")
            return

        print("  Loading layout model...")
        layout_model, layout_tokenizer = FastLanguageModel.from_pretrained(
            model_name     = layout_path,
            max_seq_length = 1024,
            load_in_4bit   = True,
            dtype          = None,
        )
        FastLanguageModel.for_inference(layout_model)
        print("  Layout model ready!")

    except Exception as e:
        print(f"  Model load failed: {e} — using rules")


# ─────────────────────────────────────
# LAYOUT DECISIONS
# ─────────────────────────────────────

def get_layout(slide_data: dict) -> dict:
    slide_type  = slide_data.get("slide_type",  "bullet_points")
    has_chart   = slide_data.get("data")         is not None
    has_bullets = bool(slide_data.get("bullets"))
    has_sub     = bool(slide_data.get("subtitle"))
    slide_num   = slide_data.get("slide_number", 1)

    if layout_model is not None:
        try:
            import torch
            prompt = f"""### Instruction:
You are a presentation layout expert. Given slide info, return layout decisions as JSON only.

### Input:
slide_type: {slide_type}, theme: corporate, slide_number: {slide_num}, total_slides: 10, has_bullets: {has_bullets}, has_chart: {has_chart}, has_subtitle: {has_sub}

### Output:
"""
            device  = "cuda" if torch.cuda.is_available() else "cpu"
            inputs  = layout_tokenizer(prompt, return_tensors="pt").to(device)
            outputs = layout_model.generate(
                **inputs, max_new_tokens=400,
                temperature=0.1, do_sample=False,
            )
            result = layout_tokenizer.decode(outputs[0], skip_special_tokens=True)
            raw    = result.split("### Output:")[-1].strip()
            if "```json" in raw:
                raw = raw.split("```json")[1].split("```")[0].strip()
            elif "```" in raw:
                raw = raw.split("```")[1].split("```")[0].strip()
            parsed = json.loads(raw)
            return parsed
        except Exception as e:
            print(f"    Model error: {e}")

    # Rule-based fallback
    rules = {
        "hero":          {"title_font_size": 52, "body_font_size": 26, "text_alignment": "center"},
        "closing":       {"title_font_size": 52, "body_font_size": 26, "text_alignment": "center"},
        "quote":         {"title_font_size": 30, "body_font_size": 22, "text_alignment": "center"},
        "bullet_points": {"title_font_size": 36, "body_font_size": 20, "text_alignment": "left"},
        "agenda":        {"title_font_size": 36, "body_font_size": 20, "text_alignment": "left"},
        "timeline":      {"title_font_size": 34, "body_font_size": 18, "text_alignment": "left"},
        "comparison":    {"title_font_size": 34, "body_font_size": 19, "text_alignment": "left"},
        "stats":         {"title_font_size": 36, "body_font_size": 18, "text_alignment": "left"},
        "chart":         {"title_font_size": 34, "body_font_size": 18, "text_alignment": "left"},
        "intro":         {"title_font_size": 38, "body_font_size": 20, "text_alignment": "left"},
        "team":          {"title_font_size": 36, "body_font_size": 18, "text_alignment": "center"},
        "image_text":    {"title_font_size": 36, "body_font_size": 20, "text_alignment": "left"},
    }
    return rules.get(slide_type, rules["bullet_points"])


# ─────────────────────────────────────
# CHART GENERATOR
# ─────────────────────────────────────

def generate_chart(data: dict, accent_color: str, slide_num: int) -> str:
    fig, ax = plt.subplots(figsize=(6, 4))
    fig.patch.set_alpha(0)

    chart_type = data.get("type", "bar_chart")
    labels     = data.get("labels", [])
    values     = data.get("values", [])

    if chart_type == "bar_chart":
        bars = ax.bar(labels, values,
                      color=accent_color, edgecolor="white", linewidth=0.5)
        ax.set_facecolor("#F8F9FA")
        if values:
            for bar, val in zip(bars, values):
                ax.text(
                    bar.get_x() + bar.get_width() / 2,
                    bar.get_height() + max(values) * 0.02,
                    str(val), ha="center", va="bottom",
                    fontsize=10, color="#333333"
                )
        if data.get("xlabel"): ax.set_xlabel(data["xlabel"], fontsize=11)
        if data.get("ylabel"): ax.set_ylabel(data["ylabel"], fontsize=11)

    elif chart_type == "line_chart":
        ax.plot(labels, values, color=accent_color,
                linewidth=2.5, marker="o", markersize=7)
        ax.fill_between(range(len(values)), values,
                        alpha=0.15, color=accent_color)
        ax.set_facecolor("#F8F9FA")
        ax.set_xticks(range(len(labels)))
        ax.set_xticklabels(labels)
        if data.get("xlabel"): ax.set_xlabel(data["xlabel"], fontsize=11)
        if data.get("ylabel"): ax.set_ylabel(data["ylabel"], fontsize=11)

    elif chart_type == "pie_chart":
        pie_colors = [accent_color, "#4CAF50", "#FF9800", "#E91E63", "#9C27B0"]
        wedges, texts, autotexts = ax.pie(
            values, labels=labels, autopct="%1.0f%%",
            colors=pie_colors[:len(values)],
            startangle=90, pctdistance=0.82
        )
        for t in autotexts:
            t.set_fontsize(11)

    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)

    ax.grid(True, axis='y', linestyle='--', alpha=0.3)
    plt.tight_layout(pad=0.5)

    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    temp_dir = os.path.join(base_dir, "temp")
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, f"chart_{slide_num}.png")
    plt.savefig(path, dpi=150, bbox_inches="tight", transparent=True)
    plt.close()
    return path


# ─────────────────────────────────────
# SMART TEXT REMOVAL — PRESERVE IMAGES
# Only clears text content, keeps shapes
# and any images/decorative elements
# ─────────────────────────────────────

def _shape_has_image(sp_elem):
    """Check if a <p:sp> element contains an image fill (blipFill)."""
    from pptx.oxml.ns import qn
    # Check for blipFill inside the shape (image fill on the shape itself)
    blip_fills = sp_elem.findall('.//' + qn('a:blipFill'))
    if blip_fills:
        return True
    # Check for blip (embedded image reference)
    blips = sp_elem.findall('.//' + qn('a:blip'))
    if blips:
        return True
    return False


def _shape_has_text(sp_elem):
    """Check if a <p:sp> element has any actual text content."""
    from pptx.oxml.ns import qn
    txBody = sp_elem.find(qn('p:txBody'))
    if txBody is None:
        return False
    for p_elem in txBody.findall(qn('a:p')):
        for r_elem in p_elem.findall(qn('a:r')):
            t_elem = r_elem.find(qn('a:t'))
            if t_elem is not None and t_elem.text and t_elem.text.strip():
                return True
    return False


def _get_shape_bounds(sp_elem):
    """Extract position and size from a shape's spPr/xfrm."""
    from pptx.oxml.ns import qn
    spPr = sp_elem.find(qn('p:spPr'))
    if spPr is None:
        return None
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        return None
    off = xfrm.find(qn('a:off'))
    ext = xfrm.find(qn('a:ext'))
    if off is None or ext is None:
        return None
    try:
        return {
            'left':   int(off.get('x', 0)),
            'top':    int(off.get('y', 0)),
            'width':  int(ext.get('cx', 0)),
            'height': int(ext.get('cy', 0)),
        }
    except (ValueError, TypeError):
        return None


def _clear_text_from_shape(sp_elem):
    """Remove all text runs from a shape's txBody, but keep the shape element."""
    from pptx.oxml.ns import qn
    txBody = sp_elem.find(qn('p:txBody'))
    if txBody is None:
        return
    # Remove all <a:p> paragraph elements
    paragraphs = txBody.findall(qn('a:p'))
    for p in paragraphs:
        txBody.remove(p)
    # Add one empty paragraph so the shape stays valid XML
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    empty_p = etree.SubElement(txBody, qn('a:p'))
    etree.SubElement(empty_p, qn('a:endParaRPr'))


def remove_text_keep_images(slide):
    """
    Smart removal: clears text content from shapes but KEEPS:
    - Shapes that contain images (blipFill)
    - The shape elements themselves (backgrounds, borders, fills)
    - Picture elements (<p:pic>)
    - Group shapes with images
    
    Returns: list of text box positions that were cleared
    """
    from pptx.oxml.ns import qn

    sp_tree = slide.shapes._spTree
    cleared_positions = []
    cleared_count = 0

    for elem in list(sp_tree):
        tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag

        if tag == 'sp':
            # Check if this shape has an image — if so, skip entirely
            if _shape_has_image(elem):
                continue

            # Check if it has text content
            if _shape_has_text(elem):
                # Record the position before clearing
                bounds = _get_shape_bounds(elem)
                if bounds and bounds['width'] > 0 and bounds['height'] > 0:
                    cleared_positions.append(bounds)

                # Clear the text but keep the shape (preserves fill, border, etc.)
                _clear_text_from_shape(elem)
                cleared_count += 1

        # Never touch picture elements, group shapes, or connectors

    return cleared_count, cleared_positions


def get_template_text_areas(cleared_positions, slide_w_emu, slide_h_emu):
    """
    From the list of cleared text positions, identify the title area and body area.
    
    Heuristic:
    - Title: the text box closest to the top with reasonable width
    - Body: the largest text box below the title
    
    Returns dict with 'title' and 'body' keys, each containing
    {left, top, width, height} in INCHES, or None if not found.
    """
    if not cleared_positions:
        return {'title': None, 'body': None}

    # Convert EMU to inches (1 inch = 914400 EMU)
    EMU_PER_INCH = 914400

    # Filter out very small shapes (decorative dots, lines etc)
    min_width = slide_w_emu * 0.15   # at least 15% of slide width
    min_height = slide_h_emu * 0.03  # at least 3% of slide height
    
    text_areas = [
        p for p in cleared_positions
        if p['width'] >= min_width and p['height'] >= min_height
    ]

    if not text_areas:
        return {'title': None, 'body': None}

    # Sort by top position (highest first)
    text_areas.sort(key=lambda p: p['top'])

    title_area = text_areas[0]  # topmost text box = title
    body_area = None

    if len(text_areas) > 1:
        # Find the largest box below the title
        candidates = text_areas[1:]
        candidates.sort(key=lambda p: p['width'] * p['height'], reverse=True)
        body_area = candidates[0]

    def to_inches(bounds):
        return {
            'left':   bounds['left']   / EMU_PER_INCH,
            'top':    bounds['top']    / EMU_PER_INCH,
            'width':  bounds['width']  / EMU_PER_INCH,
            'height': bounds['height'] / EMU_PER_INCH,
        }

    return {
        'title': to_inches(title_area),
        'body':  to_inches(body_area) if body_area else None,
    }


# ─────────────────────────────────────
# ADD CLEAN TEXTBOX
# ─────────────────────────────────────

def add_textbox(slide, left_in, top_in, width_in, height_in,
                text_content, font_size_pt,
                font_name=None,
                bold=False, color_rgb=(0, 0, 0),
                alignment=PP_ALIGN.LEFT,
                is_bullets=False,
                word_wrap=True,
                space_after_pt=10):
    """
    Adds a fresh clean textbox to the slide
    No old template formatting baggage
    """
    tb  = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in)
    )
    tf  = tb.text_frame
    tf.word_wrap = word_wrap

    if is_bullets and isinstance(text_content, list):
        for i, text in enumerate(text_content):
            if not text:
                continue
            p   = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text           = str(text)
            run.font.size      = Pt(font_size_pt)
            if font_name:
                run.font.name = font_name
            run.font.bold      = bold
            run.font.color.rgb = RGBColor(*color_rgb)
            p.alignment        = alignment
            p.space_after      = Pt(space_after_pt)
    else:
        if isinstance(text_content, list):
            text_content = "\n".join(text_content)
        p   = tf.paragraphs[0]
        run = p.add_run()
        run.text           = str(text_content)
        run.font.size      = Pt(font_size_pt)
        if font_name:
            run.font.name = font_name
        run.font.bold      = bold
        run.font.color.rgb = RGBColor(*color_rgb)
        p.alignment        = alignment

    return tb


# ─────────────────────────────────────
# DETECT BACKGROUND COLOR
# So we can choose white or dark text
# ─────────────────────────────────────

def get_text_colors_for_slide(slide, slide_type: str, registry_font_hex: str = "#0F172A"):
    """
    Returns (title_color, body_color, accent_color) as RGB tuples
    based on whether slide has a dark or light background
    """
    # Slides with full dark background image → use white text
    DARK_BG_TYPES = ["hero", "closing", "quote"]
    if slide_type in DARK_BG_TYPES:
        return (255, 255, 255), (220, 230, 245), (100, 180, 255)

    # Use registry font color for default
    fh = registry_font_hex.lstrip("#")
    font_rgb = tuple(int(fh[i:i+2], 16) for i in (0, 2, 4))
    
    # Body color slightly lighter/different if registry color is dark, otherwise same
    return font_rgb, font_rgb, (0, 120, 200)


# ─────────────────────────────────────
# BUILD SLIDE CONTENT
# Uses template text positions when available
# Falls back to proportional positioning
# ─────────────────────────────────────

def build_slide_content(slide, slide_data: dict, layout: dict,
                        accent_hex: str, font_hex: str, slide_w: float, slide_h: float,
                        template_positions: dict = None):
    """
    Adds clean textboxes and charts to the slide.
    Uses detected template text positions when available,
    otherwise falls back to proportional positioning.
    """
    from pptx.enum.shapes import MSO_SHAPE

    slide_type = slide_data.get("slide_type", "bullet_points")
    title_text = slide_data.get("title",      "") or ""
    subtitle   = slide_data.get("subtitle",   "") or ""
    bullets    = slide_data.get("bullets")    or []
    data       = slide_data.get("data")
    slide_num  = slide_data.get("slide_number", 1)

    title_pt   = layout.get("title_font_size", 36)
    body_pt    = layout.get("body_font_size",  20)
    alignment  = layout.get("text_alignment",  "left")
    align_val  = PP_ALIGN.CENTER if alignment == "center" else PP_ALIGN.LEFT

    # Accent color from hex
    ah = accent_hex.lstrip("#")
    accent_rgb = tuple(int(ah[i:i+2], 16) for i in (0, 2, 4))

    # Font Theme selection
    theme_name = slide_data.get("theme", "modern")
    font_theme = FONT_THEMES.get(theme_name, FONT_THEMES["modern"])
    title_font = font_theme["title"]
    body_font  = font_theme["body"]

    # Text colors
    title_color, body_color, _ = get_text_colors_for_slide(slide, slide_type, font_hex)

    # Margins as fractions of slide size
    margin_l = 0.05   # 5% from left
    margin_t = 0.08   # 8% from top
    margin_r = 0.05

    has_chart = data is not None

    # ── Font size adjustment (Density Heuristics) ──
    if bullets and len(bullets) > 6:
        body_pt = max(14, body_pt - 4)
    elif bullets:
        total_chars = sum(len(str(b)) for b in bullets)
        if total_chars > 500:
            body_pt = max(14, body_pt - 3)
    elif subtitle and len(str(subtitle)) > 400:
        body_pt = max(16, body_pt - 4)

    # ── Use template positions if available (otherwise fallback) ──
    tp = template_positions or {'title': None, 'body': None}

    # Content area coordinates
    content_left  = slide_w * margin_l
    content_right = slide_w * (1 - margin_r)
    content_top   = slide_h * margin_t

    if has_chart:
        text_right = slide_w * 0.48
    else:
        text_right = content_right

    text_width = text_right - content_left

    # ── Background Decorations ──
    if slide_type in ["hero", "closing", "stat"]:
        # Add a subtle background rectangle for impact
        rect_w = slide_w * 0.9
        rect_h = slide_h * 0.8
        rect_l = (slide_w - rect_w) / 2
        rect_t = (slide_h - rect_h) / 2
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(rect_l), Inches(rect_t), 
            Inches(rect_w), Inches(rect_h)
        )
        shape.fill.solid()
        if slide_type == "stat":
            shape.fill.fore_color.rgb = RGBColor(245, 247, 250) # Light grey
        else:
            shape.fill.fore_color.rgb = RGBColor(*accent_rgb)
            shape.fill.transparency = 0.8 # 80% transparent accent
        shape.line.fill.background() # No border

    # ── Title ──
    # Use template title position if detected, otherwise use defaults
    if tp.get('title') and slide_type not in ["hero", "closing", "quote"]:
        title_left   = tp['title']['left']
        title_top    = tp['title']['top']
        text_width   = tp['title']['width']
        title_height = tp['title']['height']
    else:
        title_top    = slide_h * 0.06
        title_height = slide_h * 0.16
        title_left   = content_left

        if slide_type in ["hero", "closing"]:
            title_top    = slide_h * 0.32
            title_height = slide_h * 0.22
            title_left   = slide_w * 0.08
            text_width   = slide_w * 0.84

        elif slide_type == "quote":
            title_top    = slide_h * 0.28
            title_height = slide_h * 0.35
            title_left   = slide_w * 0.08
            text_width   = slide_w * 0.84

    add_textbox(
        slide        = slide,
        left_in      = title_left,
        top_in       = title_top,
        width_in     = text_width,
        height_in    = title_height,
        text_content = title_text,
        font_size_pt = title_pt,
        font_name    = title_font,
        bold         = True,
        color_rgb    = title_color,
        alignment    = align_val,
        word_wrap    = True,
    )

    # ── Title Underline (Only for non-impact slides) ──
    if slide_type not in ["hero", "closing", "quote"]:
        line_top = title_top + title_height * 0.7 # Approximate line position
        line_width = text_width * 0.3
        line = slide.shapes.add_connector(
            1, # MSO_CONNECTOR_TYPE.STRAIGHT
            Inches(title_left), Inches(line_top),
            Inches(title_left + line_width), Inches(line_top)
        )
        line.line.color.rgb = RGBColor(*accent_rgb)
        line.line.width = Pt(2)

    # ── Body / Subtitle ──
    # Use template body position if detected
    if tp.get('body') and slide_type not in ["hero", "closing", "quote"]:
        body_top     = tp['body']['top']
        body_height  = tp['body']['height']
        content_left = tp['body']['left']
        text_width   = tp['body']['width']
    else:
        body_top    = title_top + title_height + slide_h * 0.02
        body_height = slide_h - body_top - slide_h * 0.06  # leave 6% bottom margin

    if slide_type in ["hero", "closing"] and subtitle:
        add_textbox(
            slide        = slide,
            left_in      = title_left,
            top_in       = body_top,
            width_in     = text_width,
            height_in    = slide_h * 0.10,
            text_content = subtitle,
            font_size_pt = body_pt,
            font_name    = body_font,
            bold         = False,
            color_rgb    = body_color,
            alignment    = PP_ALIGN.CENTER,
            word_wrap    = True,
        )

    elif slide_type == "quote":
        if subtitle:
            # Attribution line
            add_textbox(
                slide        = slide,
                left_in      = title_left,
                top_in       = body_top,
                width_in     = text_width,
                height_in    = slide_h * 0.10,
                text_content = subtitle,
                font_size_pt = body_pt,
                font_name    = body_font,
                bold         = False,
                color_rgb    = (accent_rgb[0], accent_rgb[1], accent_rgb[2]),
                alignment    = PP_ALIGN.CENTER,
                word_wrap    = True,
            )

    elif bullets:
        add_textbox(
            slide        = slide,
            left_in      = content_left,
            top_in       = body_top,
            width_in     = text_width,
            height_in    = body_height,
            text_content = bullets,
            font_size_pt = body_pt,
            font_name    = body_font,
            bold         = False,
            color_rgb    = body_color,
            alignment    = align_val,
            is_bullets   = True,
            word_wrap    = True,
            space_after_pt = int(slide_h * 2.5),  # proportional spacing
        )

    elif subtitle:
        add_textbox(
            slide        = slide,
            left_in      = content_left,
            top_in       = body_top,
            width_in     = text_width,
            height_in    = body_height,
            text_content = subtitle,
            font_size_pt = body_pt + 2,
            font_name    = body_font,
            bold         = False,
            color_rgb    = body_color,
            alignment    = align_val,
            word_wrap    = True,
        )

    # ── Chart ──
    if has_chart:
        chart_path = generate_chart(data, accent_hex, slide_num)

        chart_left   = slide_w * 0.50
        chart_top    = slide_h * 0.12
        chart_width  = slide_w * 0.46
        chart_height = slide_h * 0.78

        slide.shapes.add_picture(
            chart_path,
            Inches(chart_left),
            Inches(chart_top),
            Inches(chart_width),
            Inches(chart_height),
        )
        print(f"             Chart -> {data.get('type')}")

    # ── Slide Number ──
    num_txt = str(slide_num)
    add_textbox(
        slide = slide,
        left_in = slide_w - 0.6,
        top_in = slide_h - 0.5,
        width_in = 0.5,
        height_in = 0.3,
        text_content = num_txt,
        font_size_pt = 10,
        font_name = body_font,
        color_rgb = (150, 150, 150),
        alignment = PP_ALIGN.RIGHT
    )


# ─────────────────────────────────────
# SAFE SLIDE ACCESS
# Avoids the rId bug in python-pptx
# ─────────────────────────────────────

def get_slide(prs, index: int):
    """Gets slide by index safely"""
    NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldId = prs.slides._sldIdLst[index]
    rId   = sldId.get(f'{{{NS}}}id')
    return prs.slides.part.related_slide(rId)


def get_slide_count(prs) -> int:
    return len(prs.slides._sldIdLst)


# ─────────────────────────────────────
# DUPLICATE A SLIDE
# ─────────────────────────────────────

def duplicate_slide(prs, index: int):
    """Duplicates slide at index, returns new slide"""
    NS     = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldId  = prs.slides._sldIdLst[index]
    rId    = sldId.get(f'{{{NS}}}id')
    src    = prs.slides.part.related_slide(rId)

    xml_copy  = copy.deepcopy(src._element)
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])
    new_slide.part._element = xml_copy
    return prs.slides[-1]


# ─────────────────────────────────────
# REMOVE EXTRA SLIDES
# ─────────────────────────────────────

def remove_last_slide(prs):
    """Removes the last slide safely"""
    NS    = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldId = prs.slides._sldIdLst[-1]
    rId   = sldId.get(f'{{{NS}}}id')
    try:
        prs.part.drop_rel(rId)
    except:
        pass
    del prs.slides._sldIdLst[-1]


# ─────────────────────────────────────
# MAIN RUN FUNCTION
# ─────────────────────────────────────

def run(
    content_json:         dict,
    output_path:          str,
    template_id:          str  = "modern_blue",
    custom_template_path: str  = None
):
    import random
    from template_registry import TEMPLATES

    load_models()

    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

    # ── Pick template ──
    accent_color = "#00A8E8"
    font_color   = "#0F172A"
    template_path = None

    # Handle "random" or None template_id
    if (not template_id or template_id == "random") and TEMPLATES:
        template_id = random.choice(list(TEMPLATES.keys()))
        print(f"Randomly selected template: {template_id}")

    if custom_template_path and os.path.exists(custom_template_path):
        template_path = custom_template_path
        accent_color  = "#0066CC"
        print(f"Using uploaded template: {custom_template_path}")

    elif template_id and template_id in TEMPLATES:
        template_path = os.path.join(base_dir, TEMPLATES[template_id]["file"])
        accent_color  = TEMPLATES[template_id].get("accent_color", "#00A8E8")
        font_color    = TEMPLATES[template_id].get("font_color", "#0F172A")
        print(f"Using template: {TEMPLATES[template_id]['name']}")

    else:
        # Check if there are any "ppt template XX" keys
        fallback_templates = [k for k in TEMPLATES.keys() if "ppt template" in k]
        if fallback_templates:
            template_id = random.choice(fallback_templates)
            template_path = os.path.join(base_dir, TEMPLATES[template_id]["file"])
            accent_color = TEMPLATES[template_id].get("accent_color", "#00A8E8")
            font_color   = TEMPLATES[template_id].get("font_color", "#0F172A")
            print(f"Using fallback user template: {template_id}")
        else:
            # Fixed absolute fallback
            template_path = os.path.join(base_dir, "templates/Bold Text and Color Morph Product Launch Presentation.pptx")
            print("Using hardcoded fallback template")

    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")

    # ── Open template ──
    prs = Presentation(template_path)

    slide_w = prs.slide_width.inches
    slide_h = prs.slide_height.inches
    print(f"Slide size: {slide_w:.2f} x {slide_h:.2f} inches")

    original_count = get_slide_count(prs)
    slides_content = content_json.get("slides", [])

    print(f"Template slides: {original_count}")
    print(f"Content slides:  {len(slides_content)}\n")

    # ── Process each content slide ──
    for i, slide_data in enumerate(slides_content):
        title  = slide_data.get("title", "")[:45]
        stype  = slide_data.get("slide_type", "")
        print(f"  Slide {i+1:02d}/{len(slides_content)} [{stype:<15}] {title}")

        # Get layout decisions from model
        layout = get_layout(slide_data)

        # ── Get or create slide ──
        if i < original_count:
            slide = get_slide(prs, i)
        else:
            # Duplicate a template slide for extra content slides
            dup_from = min(1, original_count - 1)
            slide    = duplicate_slide(prs, dup_from)
            print(f"             Duplicated template slide {dup_from}")

        # ── STEP 1: Smart text removal — preserve images ──
        cleared_count, cleared_positions = remove_text_keep_images(slide)
        print(f"             Cleared text from {cleared_count} shapes (images preserved)")

        # ── STEP 1b: Detect template text positions ──
        slide_w_emu = prs.slide_width
        slide_h_emu = prs.slide_height
        template_positions = get_template_text_areas(
            cleared_positions, slide_w_emu, slide_h_emu
        )
        if template_positions.get('title'):
            print(f"             Using template title position")
        if template_positions.get('body'):
            print(f"             Using template body position")

        # ── STEP 2: Add our clean content using template positions ──
        build_slide_content(
            slide              = slide,
            slide_data         = slide_data,
            layout             = layout,
            accent_hex         = accent_color,
            font_hex           = font_color,
            slide_w            = slide_w,
            slide_h            = slide_h,
            template_positions = template_positions,
        )
        print(f"             Done")

    # ── Remove extra template slides ──
    current = get_slide_count(prs)
    extra   = current - len(slides_content)
    if extra > 0:
        print(f"\nRemoving {extra} extra template slides...")
        for _ in range(extra):
            try:
                remove_last_slide(prs)
            except Exception as e:
                print(f"  Remove failed: {e}")
                break

    # ── Save ──
    out_dir = os.path.dirname(output_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)

    prs.save(output_path)
    final_count = get_slide_count(prs)
    print(f"\n{'='*45}")
    print(f"Saved:  {output_path}")
    print(f"Slides: {final_count}")
    print(f"{'='*45}")


# ─────────────────────────────────────
# DIRECT TEST
# python agents/visual_agent.py content.json
# ─────────────────────────────────────

if __name__ == "__main__":
    import sys
    sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

    if len(sys.argv) > 1:
        with open(sys.argv[1]) as f:
            content = json.load(f)
        run(content, "test_output.pptx")
        print("Done! Open test_output.pptx")
    else:
        print("Usage: python agents/visual_agent.py <content.json>")